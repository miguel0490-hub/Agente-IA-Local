"""
document_parser.py — Router Universal de Archivos
==================================================
Extrae contenido legible de CUALQUIER archivo para enviarlo al LLM.

Arquitectura (Strategy Pattern):
  1. Si la extensión tiene un extractor dedicado → úsalo.
  2. Si no → intenta lectura como texto plano (UTF-8, con fallback a latin-1).
  3. Si el archivo es binario ininteligible → devuelve un mensaje de error
     manejado, nunca lanza una excepción sin atrapar.

Regla de Chesterton's Fence:
  Se conservan TODOS los parsers anteriores (_parse_pdf, _parse_docx, etc.)
  ya que son consumidos activamente por el flujo de chat existente.
  Solo se expande el comportamiento de `extraer_texto_archivo` hacia abajo.
"""

import os
import io
import base64
from pathlib import Path

from src.core.i18n import t

# ─── Extensiones clasificadas como imágenes ────────────────────────────────────
_IMAGEN_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp', '.gif', '.bmp', '.tiff', '.tif', '.ico', '.svg'}

# ─── Extensiones clasificadas como vídeo (pasan por ruta de Gemini) ────────────
_VIDEO_EXTENSIONS = {'.mp4', '.mov', '.avi', '.mkv', '.webm', '.flv', '.wmv'}

# ─── Binarios conocidos que NUNCA tienen texto legible ─────────────────────────
_BINARY_EXTENSIONS = {
    '.exe', '.dll', '.so', '.dylib', '.bin', '.dat', '.db', '.sqlite',
    '.zip', '.tar', '.gz', '.rar', '.7z', '.bz2',
    '.pyc', '.pyo', '.class', '.jar', '.war',
    '.mp3', '.wav', '.ogg', '.flac', '.m4a', '.aac',
    '.ttf', '.otf', '.woff', '.woff2',
}

_AUDIO_EXTENSIONS = {".mp3", ".wav", ".ogg", ".flac", ".m4a", ".aac"}

# ─── Parsers dedicados (Patrón Strategy) ──────────────────────────────────────


def _parse_pdf(file_obj) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_obj)
    paginas = [page.extract_text() for page in reader.pages if page.extract_text()]
    if not paginas:
        return t("doc_pdf_no_text")
    return "\n".join(paginas)


def _parse_docx(file_obj) -> str:
    from docx import Document
    doc = Document(file_obj)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])


def _parse_odf(file_obj) -> str:
    from odf.opendocument import load as odf_load
    from odf.teletype import extractText as odf_extract_text
    doc = odf_load(file_obj)
    return odf_extract_text(doc)


def _parse_excel(file_obj, is_ods: bool = False) -> str:
    import pandas as pd
    motor = 'odf' if is_ods else None
    df = pd.read_excel(file_obj, engine=motor)
    return t("doc_spreadsheet_prefix") + df.to_string()


def _parse_csv(file_obj) -> str:
    import pandas as pd
    # Intentar con UTF-8 primero, fallback a latin-1
    try:
        df = pd.read_csv(file_obj)
    except UnicodeDecodeError:
        file_obj.seek(0)
        df = pd.read_csv(file_obj, encoding='latin-1')
    return t("doc_csv_prefix") + df.to_string()


def _parse_pptx(file_obj) -> str:
    from pptx import Presentation
    prs = Presentation(file_obj)
    texto = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texto.append(shape.text)
    return "\n".join(texto)


def _parse_text(file_obj) -> str:
    """Lee cualquier archivo de texto. Intenta UTF-8, luego latin-1."""
    raw = file_obj.read()
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("latin-1", errors="replace")


def _parse_json(file_obj) -> str:
    import json
    try:
        data = json.load(file_obj)
        return t("doc_json_header") + json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError as e:
        file_obj.seek(0)
        body = _parse_text(file_obj)
        return t("doc_json_malformed", body=body, error=e)


def _parse_image_as_description(file_obj) -> str:
    """
    Para imágenes: devuelve metadata básica + Base64 del contenido.
    El LLM de visión (Gemini) ya tiene su propia ruta en app.py para procesar
    imágenes como objeto PIL. Este parser solo actúa como fallback de contexto.
    """
    raw = file_obj.read()
    nombre = getattr(file_obj, 'name', t("doc_unnamed_image"))
    ext = Path(nombre).suffix.lower().lstrip('.')
    b64 = base64.b64encode(raw).decode('utf-8')
    size_kb = len(raw) / 1024
    preview = b64[:200]
    return t(
        "doc_image_attachment",
        name=nombre,
        size_kb=f"{size_kb:.1f}",
        ext_label=ext.upper(),
        mime_ext=ext,
        preview=preview,
    )


# ─── Mapa de estrategias (extensión → parser) ─────────────────────────────────
_EXTRACTORS: dict = {
    # Documentos ricos
    '.pdf':  _parse_pdf,
    '.docx': _parse_docx,
    '.odt':  _parse_odf,
    '.odp':  _parse_odf,
    '.pptx': _parse_pptx,
    # Hojas de cálculo
    '.xlsx': _parse_excel,
    '.xls':  _parse_excel,
    '.ods':  lambda f: _parse_excel(f, is_ods=True),
    '.csv':  _parse_csv,
    # Texto y código (todos pasan por _parse_text)
    '.txt':  _parse_text,
    '.md':   _parse_text,
    '.py':   _parse_text,
    '.js':   _parse_text,
    '.ts':   _parse_text,
    '.jsx':  _parse_text,
    '.tsx':  _parse_text,
    '.html': _parse_text,
    '.htm':  _parse_text,
    '.css':  _parse_text,
    '.scss': _parse_text,
    '.xml':  _parse_text,
    '.yaml': _parse_text,
    '.yml':  _parse_text,
    '.toml': _parse_text,
    '.ini':  _parse_text,
    '.env':  _parse_text,
    '.sh':   _parse_text,
    '.bat':  _parse_text,
    '.sql':  _parse_text,
    '.rs':   _parse_text,
    '.go':   _parse_text,
    '.java': _parse_text,
    '.cpp':  _parse_text,
    '.c':    _parse_text,
    '.h':    _parse_text,
    '.php':  _parse_text,
    '.rb':   _parse_text,
    '.kt':   _parse_text,
    '.swift': _parse_text,
    '.r':    _parse_text,
    '.log':  _parse_text,
    '.conf': _parse_text,
    '.cfg':  _parse_text,
    # JSON (parser dedicado con pretty-print)
    '.json': _parse_json,
    '.jsonl': _parse_text,
    # Imágenes (fallback informativo — la ruta Gemini Vision es la principal)
    **{ext: _parse_image_as_description for ext in _IMAGEN_EXTENSIONS},
}


# ─── Fallback Universal ────────────────────────────────────────────────────────

def _fallback_universal(file_obj, nombre: str) -> str:
    """
    Último recurso para archivos sin extractor dedicado.
    Intenta leer como texto UTF-8 → latin-1.
    Si es binario ininteligible, devuelve un mensaje de error manejado.
    """
    ext = Path(nombre).suffix.lower()

    # Cortocircuito rápido para binarios conocidos
    if ext in _BINARY_EXTENSIONS:
        if ext in _AUDIO_EXTENSIONS:
            return t("doc_audio_binary", name=nombre, ext=ext)
        return t("doc_binary_generic", name=nombre, ext=ext)

    # Intentar lectura de texto para extensiones desconocidas
    try:
        raw = file_obj.read()
        # Heurística: si más del 30% de los primeros 512 bytes son nulos o no imprimibles → binario
        muestra = raw[:512]
        bytes_no_imprimibles = sum(1 for b in muestra if b < 9 or (14 <= b <= 31) or b == 127)
        if len(muestra) > 0 and (bytes_no_imprimibles / len(muestra)) > 0.30:
            ext_display = ext or "?"
            return t("doc_binary_heuristic", name=nombre, ext=ext_display)
        # Es texto — decodificar
        try:
            return raw.decode("utf-8")
        except UnicodeDecodeError:
            return raw.decode("latin-1", errors="replace")

    except Exception as e:
        return t("doc_read_error", name=nombre, error=e)


# ─── Función Pública Principal ─────────────────────────────────────────────────

def extraer_texto_archivo(file_obj) -> str:
    """
    Router Universal de Archivos.

    Recibe cualquier objeto de archivo (Streamlit UploadedFile o file-like)
    y devuelve SIEMPRE un string con el contenido legible o un mensaje
    de error descriptivo. NUNCA devuelve None.

    Flujo de decisión:
      1. Extractor dedicado por extensión → usa el parser específico.
      2. Sin extractor → fallback universal (heurística de texto/binario).
      3. Cualquier excepción interna → capturada y convertida a mensaje de error.
    """
    nombre = getattr(file_obj, 'name', t("doc_unnamed_file"))
    ext = Path(nombre.lower()).suffix

    # Early return: extensiones de vídeo — se procesan por ruta separada en app.py
    if ext in _VIDEO_EXTENSIONS:
        return t("doc_video_placeholder", name=nombre)

    extractor = _EXTRACTORS.get(ext)

    if extractor:
        try:
            texto_extraido = extractor(file_obj)
        except Exception as e:
            file_obj.seek(0)
            fb = _fallback_universal(file_obj, nombre)
            texto_extraido = t("doc_parse_error", name=nombre, ext=ext, error=e, fallback=fb)
    else:
        # Sin extractor dedicado → Fallback Universal
        texto_extraido = _fallback_universal(file_obj, nombre)

    # Evaluación para RAG (Archivos muy grandes > 5000 palabras)
    palabras = len(texto_extraido.split())
    if palabras > 5000 and not texto_extraido.startswith("⛔"):
        from src.services.task_queue import enqueue_rag_indexing
        from src.services.rag_service import RAGService

        job_id = enqueue_rag_indexing(nombre, texto_extraido)
        if job_id:
            return t("doc_rag_enqueued", name=nombre, words=palabras, job_id=job_id)

        rag = RAGService()
        chunks = rag.index_document(nombre, texto_extraido)
        return t("doc_rag_indexed", name=nombre, words=palabras, chunks=chunks)

    return texto_extraido
